"""Generate one-slide PPTX summary with scatter visualization."""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


HIGHLIGHT_ORDER = ["top_rated", "lowest_rated", "highest_volume"]
HIGHLIGHT_META = {
    "top_rated": {
        "legend": "Top Rated",
        "short": "Top",
        "color": "#2E8B57",
        "marker": "*",
        "size": 190,
    },
    "lowest_rated": {
        "legend": "Lowest Rated",
        "short": "Low",
        "color": "#C0392B",
        "marker": "X",
        "size": 130,
    },
    "highest_volume": {
        "legend": "Highest Volume",
        "short": "Volume",
        "color": "#E67E22",
        "marker": "D",
        "size": 125,
    },
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--metrics-csv", default="outputs/final/location_metrics.csv")
    parser.add_argument("--enriched-csv", default="outputs/final/reviews_enriched.csv")
    parser.add_argument("--output-pptx", default="outputs/final/assessment_insights.pptx")
    parser.add_argument("--title", default="la Madeleine Location Review Insights")
    parser.add_argument("--plot-path", default="outputs/plots/review_scatter.png")
    return parser.parse_args()


def _normalize_metrics_for_analysis(metrics: pd.DataFrame) -> pd.DataFrame:
    metrics_plot = metrics.copy()
    metrics_plot["review_count"] = pd.to_numeric(metrics_plot["review_count"], errors="coerce")
    metrics_plot["avg_rating"] = pd.to_numeric(metrics_plot["avg_rating"], errors="coerce")
    metrics_plot = metrics_plot.dropna(subset=["review_count", "avg_rating"])
    return metrics_plot


def format_location_label(row: pd.Series) -> str:
    name = str(row.get("locationName", "Unknown")).strip()
    city = str(row.get("city", "")).strip()
    state = str(row.get("state", "")).strip()
    if city and state:
        return f"{name} ({city}, {state})"
    if city:
        return f"{name} ({city})"
    if state:
        return f"{name} ({state})"
    return name


def select_highlight_points(metrics: pd.DataFrame) -> dict[str, Any]:
    metrics_clean = _normalize_metrics_for_analysis(metrics)
    if metrics_clean.empty:
        return {
            "threshold": 50,
            "metrics_clean": metrics_clean,
            "stable_subset": metrics_clean,
            "top_rated": None,
            "lowest_rated": None,
            "highest_volume": None,
        }

    threshold = max(int(metrics_clean["review_count"].median()), 50)
    stable_subset = metrics_clean.loc[metrics_clean["review_count"] >= threshold]
    if stable_subset.empty:
        stable_subset = metrics_clean.copy()

    top_rated = stable_subset.sort_values(by=["avg_rating", "review_count"], ascending=[False, False]).iloc[0]
    lowest_rated = stable_subset.sort_values(by=["avg_rating", "review_count"], ascending=[True, False]).iloc[0]
    highest_volume = metrics_clean.sort_values(by=["review_count", "avg_rating"], ascending=[False, False]).iloc[0]

    return {
        "threshold": threshold,
        "metrics_clean": metrics_clean,
        "stable_subset": stable_subset,
        "top_rated": top_rated,
        "lowest_rated": lowest_rated,
        "highest_volume": highest_volume,
    }


def _build_point_key(row: pd.Series) -> str:
    store_id = str(row.get("storeID", "")).strip()
    if store_id:
        return store_id
    return format_location_label(row)


def create_scatter_plot(metrics: pd.DataFrame, highlights: dict[str, Any], plot_path: str) -> Path:
    output_path = Path(plot_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    metrics_plot = highlights["metrics_clean"]

    fig, ax = plt.subplots(figsize=(7.5, 4.2), dpi=160)
    ax.scatter(
        metrics_plot["review_count"],
        metrics_plot["avg_rating"],
        s=45,
        alpha=0.55,
        color="#5B8FB9",
        edgecolor="white",
        linewidth=0.5,
    )

    if not metrics_plot.empty and highlights["top_rated"] is not None:
        unique_points: dict[str, dict[str, Any]] = {}
        for category in HIGHLIGHT_ORDER:
            row = highlights[category]
            key = _build_point_key(row)
            if key not in unique_points:
                unique_points[key] = {
                    "x": float(row["review_count"]),
                    "y": float(row["avg_rating"]),
                    "row": row,
                    "categories": [],
                }
            unique_points[key]["categories"].append(category)

            meta = HIGHLIGHT_META[category]
            ax.scatter(
                [float(row["review_count"])],
                [float(row["avg_rating"])],
                s=meta["size"],
                marker=meta["marker"],
                color=meta["color"],
                edgecolor="black",
                linewidth=0.6,
                zorder=4,
            )

        offsets = [(12, 16), (14, -18), (-125, 16), (-125, -18), (30, 22), (30, -24)]
        for idx, point in enumerate(unique_points.values()):
            categories = point["categories"]
            short_tags = "/".join(HIGHLIGHT_META[c]["short"] for c in categories)
            label = f"{short_tags}: {format_location_label(point['row'])}"
            dx, dy = offsets[idx % len(offsets)]
            ax.annotate(
                label,
                xy=(point["x"], point["y"]),
                xytext=(dx, dy),
                textcoords="offset points",
                fontsize=8,
                bbox={"boxstyle": "round,pad=0.2", "fc": "white", "ec": "#666666", "alpha": 0.9},
                arrowprops={"arrowstyle": "->", "lw": 0.6, "color": "#555555"},
                zorder=5,
            )

        legend_handles = []
        for category in HIGHLIGHT_ORDER:
            meta = HIGHLIGHT_META[category]
            handle = ax.scatter(
                [],
                [],
                s=meta["size"] * 0.55,
                marker=meta["marker"],
                color=meta["color"],
                edgecolor="black",
                linewidth=0.6,
                label=meta["legend"],
            )
            legend_handles.append(handle)
        ax.legend(handles=legend_handles, loc="lower right", fontsize=8, framealpha=0.95)

    ax.set_title("Location Review Volume vs. Average Rating", fontsize=11)
    ax.set_xlabel("Review Count", fontsize=10)
    ax.set_ylabel("Average Rating", fontsize=10)
    ax.grid(alpha=0.25, linestyle="--", linewidth=0.5)

    fig.tight_layout()
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)
    return output_path


def generate_insight_bullets(highlights: dict[str, Any], enriched: pd.DataFrame) -> list[str]:
    if highlights["metrics_clean"].empty or highlights["top_rated"] is None:
        return ["No matched location metrics were available for automated insights."]

    top_rated = highlights["top_rated"]
    lowest_rated = highlights["lowest_rated"]
    highest_volume = highlights["highest_volume"]
    threshold = int(highlights["threshold"])
    stable_subset = highlights["stable_subset"]

    matched_mask = enriched.get("matched_storeID", pd.Series(dtype=str)).fillna("") != ""
    matched_count = int(matched_mask.sum())
    total_count = int(len(enriched))
    coverage = (matched_count / total_count * 100) if total_count else 0.0

    rating_spread = float(stable_subset["avg_rating"].max() - stable_subset["avg_rating"].min())

    bullets = [
        (
            f"Top rated location (>= {threshold} reviews): {format_location_label(top_rated)} "
            f"with avg rating {float(top_rated['avg_rating']):.2f} across {int(top_rated['review_count'])} reviews."
        ),
        (
            f"Lowest rated location (>= {threshold} reviews): {format_location_label(lowest_rated)} "
            f"at {float(lowest_rated['avg_rating']):.2f}; indicates service quality variance."
        ),
        (
            f"Highest volume location: {format_location_label(highest_volume)} with "
            f"{int(highest_volume['review_count'])} reviews and avg rating {float(highest_volume['avg_rating']):.2f}."
        ),
        f"Review-to-location match coverage is {coverage:.1f}% ({matched_count}/{total_count}).",
        f"Rating spread across stable locations is {rating_spread:.2f} points, useful for geo-prioritized deep dives.",
    ]
    return bullets


def build_slide(title: str, plot_path: str, output_pptx: str, bullets: list[str]) -> Path:
    output_path = Path(output_pptx)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.4), Inches(0.7))
    title_tf = title_box.text_frame
    title_run = title_tf.paragraphs[0].add_run()
    title_run.text = title
    title_run.font.size = Pt(28)
    title_run.font.bold = True

    slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.0), width=Inches(7.4), height=Inches(4.6))

    insights_box = slide.shapes.add_textbox(Inches(8.1), Inches(1.0), Inches(5.0), Inches(4.6))
    insights_tf = insights_box.text_frame
    insights_tf.word_wrap = True

    p0 = insights_tf.paragraphs[0]
    p0.text = "Key Highlights"
    p0.font.size = Pt(20)
    p0.font.bold = True

    for bullet in bullets:
        p = insights_tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(12)

    presentation.save(output_path)
    return output_path


def main() -> None:
    args = parse_args()
    metrics = pd.read_csv(args.metrics_csv)
    enriched = pd.read_csv(args.enriched_csv)

    highlights = select_highlight_points(metrics)
    plot_path = create_scatter_plot(metrics, highlights, args.plot_path)
    bullets = generate_insight_bullets(highlights, enriched)
    pptx_path = build_slide(args.title, str(plot_path), args.output_pptx, bullets)

    print(f"Created slide: {pptx_path}")


if __name__ == "__main__":
    main()
